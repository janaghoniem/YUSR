# ============================================================================
# Module Selection Layer
# ============================================================================
"""
Smart module selector that intercepts tasks and recommends the best library
before the LLM starts writing code.

This prevents the LLM from defaulting to pyautogui when better, more reliable
libraries are available (e.g., python-docx for Word, openpyxl for Excel).
"""

from dataclasses import dataclass
from typing import Optional, List
import logging
import os

logger = logging.getLogger(__name__)


# ============================================================================
# Common Data Extraction Rules (DRY Principle)
# ============================================================================

CRITICAL_DATA_EXTRACTION_RULES = """
================================================================================
CRITICAL - DATA EXTRACTION AND OUTPUT RULES (NON-NEGOTIABLE)
================================================================================

When your code performs data extraction, copying, reading, or retrieval:

**YOU MUST OUTPUT THE ACTUAL DATA BEFORE ANY STATUS MESSAGE**

This is CRITICAL for multi-step workflows where subsequent tasks depend on your output.

✅ CORRECT PATTERN:
```python
# Extract/copy/read the data
content = pyperclip.paste()  # or file.read(), or extracted_data, etc.

# OUTPUT THE ACTUAL DATA FIRST
print(content)

# THEN print success indicator
print("EXECUTION_SUCCESS")
```

❌ WRONG PATTERN (DATA IS LOST):
```python
content = pyperclip.paste()
if content:
    print("EXECUTION_SUCCESS")  # ← Only status message, actual data is lost!
```

**Why this matters:**
- The next task in the pipeline receives your printed output as its input
- If you only print "EXECUTION_SUCCESS", the next task receives an empty string
- The entire workflow fails silently
- Always output data BEFORE status messages

**FILE TYPE SPECIFIC EXTRACTION METHODS:**

PDF FILES - Use PyPDF2 for accurate text extraction:
```python
from PyPDF2 import PdfReader
import os

try:
    pdf_path = "path/to/file.pdf"  # Use actual file path

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"File not found: {{pdf_path}}")

    # Extract text from PDF with None handling
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:  # Handle None and empty strings
            text += page_text

    # Check if we extracted any content
    if not text.strip():
        text = "[No text content found in PDF]"

    # OUTPUT THE ACTUAL PDF CONTENT FIRST
    print(text)

    # THEN indicate success
    print("EXECUTION_SUCCESS")
except Exception as e:
    print(f"Exception: {{e}}")
```

COPYING FILE CONTENT (NON-PDF):
```python
import pyautogui
import pyperclip
import time

try:
    # Select all and copy
    pyautogui.hotkey('ctrl', 'a')
    time.sleep(0.2)
    pyautogui.hotkey('ctrl', 'c')
    time.sleep(0.2)

    # Get the copied content
    content = pyperclip.paste()

    # OUTPUT THE ACTUAL CONTENT FIRST
    print(content)

    # THEN indicate success
    print("EXECUTION_SUCCESS")
except Exception as e:
    print(f"Exception: {{e}}")
```

READING A FILE:
```python
import os

try:
    filepath = "D:/Downloads/file.txt"

    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {{filepath}}")

    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # OUTPUT THE FILE CONTENT FIRST
    print(content)

    # THEN indicate success
    print("EXECUTION_SUCCESS")
except Exception as e:
    print(f"Exception: {{e}}")
```

EXTRACTING TEXT FROM UI:
```python
import pyautogui
import pyperclip
import time

try:
    # Select text
    pyautogui.hotkey('ctrl', 'a')
    time.sleep(0.1)

    # Copy to clipboard
    pyautogui.hotkey('ctrl', 'c')
    time.sleep(0.2)

    # Extract from clipboard
    extracted_text = pyperclip.paste()

    # OUTPUT THE EXTRACTED TEXT FIRST
    print(extracted_text)

    # THEN indicate success
    print("EXECUTION_SUCCESS")
except Exception as e:
    print(f"Exception: {{e}}")
```

WEB SCRAPING / DATA EXTRACTION:
```python
# Whatever extraction method you use...
extracted_data = element.get_text()

# OUTPUT THE EXTRACTED DATA FIRST
print(extracted_data)

# THEN indicate success
print("EXECUTION_SUCCESS")
```

**SMART FILE TYPE DETECTION:**
- **PDF files (.pdf)**: Use PyPDF2 for text extraction
- **Office files (.docx, .xlsx, .pptx)**: Use appropriate libraries (python-docx, openpyxl, python-pptx)
- **Text files (.txt, .py, .json, etc.)**: Use file reading with UTF-8 encoding
- **Other files**: Use clipboard copying as fallback

**Remember:**
- Data output FIRST
- Status message SECOND
- Choose the RIGHT extraction method based on file type
- This applies to ANY task that extracts, copies, reads, or retrieves information
- Failure to follow this pattern breaks the entire workflow
"""


# ============================================================================
# Dynamic Folder Configuration
# ============================================================================

def get_agent_folder(subfolder: str) -> str:
    """
    Get dynamic agent folder path that works across different systems

    Args:
        subfolder: The specific subfolder ('docs', 'excel', 'ppts')

    Returns:
        Full path to the agent subfolder
    """
    import os

    # Try multiple common desktop locations - user's preference first
    possible_paths = [
        "D:/OneDrive/Desktop/agent",  # User's preferred path - FIRST PRIORITY
        os.path.expanduser("~/OneDrive/Desktop/agent"),
        os.path.expanduser("~/Desktop/agent"),
        os.path.expanduser("~/Documents/agent"),
    ]

    for base_path in possible_paths:
        if os.path.exists(os.path.dirname(base_path)) or os.path.exists(base_path):
            agent_folder = os.path.join(base_path, subfolder)
            os.makedirs(agent_folder, exist_ok=True)
            return agent_folder.replace("\\", "\\\\")

    # Final fallback - create in user's home directory
    fallback_path = os.path.join(os.path.expanduser("~"), "agent", subfolder)
    os.makedirs(fallback_path, exist_ok=True)
    return fallback_path.replace("\\", "\\\\")


# ============================================================================
# Module Definitions
# ============================================================================

@dataclass
class ModuleGuidance:
    """Guidance for the LLM about which module to use"""
    module_name: str
    library_name: str
    library_import: str
    keywords: List[str]  # Task keywords that trigger this module
    guidance: str  # Instruction block to inject into prompt

    def get_override_block(self) -> str:
        """Get the formatted override block for the prompt"""
        return f"""[MODULE OVERRIDE: {self.module_name}]
Library: {self.library_name}
Import: {self.library_import}

{self.guidance}
[/MODULE OVERRIDE]"""


# Pre-defined module overrides
MODULES = {
    "word": ModuleGuidance(
        module_name="Word",
        library_name="python-docx",
        library_import="from docx import Document",
        keywords=["word", "document", "docx", "page", "paragraph", "heading", "table", "text formatting"],
        guidance=f"""
Work with .docx files using python-docx. Identify TASK TYPE first, then follow ONLY that branch.

{CRITICAL_DATA_EXTRACTION_RULES}

SETUP:
import os, glob
from docx import Document
from datetime import datetime

# Dynamic folder detection - user's preferred path first
possible_paths = [
    "D:/OneDrive/Desktop/agent/docs",  # User's preferred path - FIRST PRIORITY
    os.path.expanduser("~/OneDrive/Desktop/agent/docs"),
    os.path.expanduser("~/Desktop/agent/docs"),
    os.path.expanduser("~/Documents/agent/docs"),
]

folder = None
for path in possible_paths:
    if os.path.exists(os.path.dirname(path)) or os.path.exists(path):
        folder = path
        break

# Final fallback
if not folder:
    folder = os.path.join(os.path.expanduser("~"), "agent", "docs")

os.makedirs(folder, exist_ok=True)
files = glob.glob(os.path.join(folder, '*.docx'))
latest = max(files, key=os.path.getctime) if files else None
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
# Read [ACTIVE FILE] from prompt if present:
active_file = '<path from [ACTIVE FILE: ...]>' if '[ACTIVE FILE:' in PROMPT else latest

TASK TYPE — match ONE, execute ONLY that branch:

[1] LAUNCH APP — task says "open word application", "open microsoft word", "launch word", "start word"
    AND does NOT mention a specific document/file to edit:
    os.system('start winword')
    print("EXECUTION_SUCCESS")
    # Done. No file operations.

[2] OPEN FILE — task mentions a specific filename or path (e.g. "open report.docx", "open existing file named X"):
    target = 'exact_filename_from_task.docx'  # extract from task description
    if os.path.isabs(target) and os.path.exists(target):
        found_path = target
    else:
        search_roots = [
            os.path.expanduser('~\\Desktop'),
            os.path.expanduser('~\\OneDrive\\Desktop'),
            os.path.expanduser('~\\Documents'),
            os.path.expanduser('~\\Downloads'),
            'D:\\OneDrive\\Desktop',
            folder,
        ]
        found_path = None
        for root in search_roots:
            if not os.path.exists(root):
                continue
            matches = glob.glob(os.path.join(root, '**', target), recursive=True)
            if matches:
                found_path = matches[0]
                break
    if not found_path:
        raise FileNotFoundError(f"Could not find file: {{target}}")
    os.startfile(found_path)
    print(f"[FILE]: {{found_path}}")
    print("EXECUTION_SUCCESS")
    # Done.

[3] CREATE — task says "create/new/make/generate" + document/file:
    save_path = os.path.join(folder, f'descriptive_name_{{ts}}.docx')
    doc = Document()
    doc.save(save_path)
    print(f"[FILE]: {{save_path}}")
    print("EXECUTION_SUCCESS")
    # Done. Do NOT open the file — keep it unlocked for subsequent edit tasks.

[4] SAVE / CONFIRM — task says "save", "press ok", "press save", "click ok", "click save",
    "confirm saving", "press enter to confirm":
    check_path = active_file
    if check_path and os.path.exists(check_path):
        os.startfile(check_path)  # Open in Word now that all edits are done
        print(f"[FILE]: {{check_path}}")
    print("EXECUTION_SUCCESS")
    # Done. NO doc.save() call.

[5] DEFAULT — task says write/type/add/modify/edit/set/insert/format/heading/paragraph:
    doc = Document(active_file) if active_file and os.path.exists(active_file) else Document()
    save_path = active_file if active_file and os.path.exists(active_file) else os.path.join(folder, f'document_{{ts}}.docx')
    # ... apply changes using python-docx ...
    try:
        doc.save(save_path)
    except PermissionError:
        save_path = os.path.splitext(save_path)[0] + '_v2.docx'
        doc.save(save_path)
    print(f"[FILE]: {{save_path}}")
    # Automatically open the file after editing operations are complete
    os.startfile(save_path)
    print("EXECUTION_SUCCESS")

Key python-docx patterns:
- doc.add_heading('Title', level=0)   # document title
- doc.add_heading('Section', level=1) # section heading
- doc.add_paragraph('text')           # body paragraph
- p.add_run('bold').bold = True

AVOID: subprocess, pyautogui, Word UI, mouse/keyboard events
"""
    ),
    "excel": ModuleGuidance(
        module_name="Excel",
        library_name="openpyxl",
        library_import="from openpyxl import Workbook",
        keywords=["excel", "spreadsheet", "xlsx", "sheet", "cell", "row", "column", "formula", "data table"],
        guidance=f"""
Work with .xlsx files using openpyxl. Identify TASK TYPE first, then follow ONLY that branch.

{CRITICAL_DATA_EXTRACTION_RULES}

SETUP:
import os, glob
from openpyxl import Workbook, load_workbook
from datetime import datetime

# Dynamic folder detection - user's preferred path first
possible_paths = [
    "D:/OneDrive/Desktop/agent/excel",  # User's preferred path - FIRST PRIORITY
    os.path.expanduser("~/OneDrive/Desktop/agent/excel"),
    os.path.expanduser("~/Desktop/agent/excel"),
    os.path.expanduser("~/Documents/agent/excel"),
]

folder = None
for path in possible_paths:
    if os.path.exists(os.path.dirname(path)) or os.path.exists(path):
        folder = path
        break

# Final fallback
if not folder:
    folder = os.path.join(os.path.expanduser("~"), "agent", "excel")

os.makedirs(folder, exist_ok=True)
files = glob.glob(os.path.join(folder, '*.xlsx'))
latest = max(files, key=os.path.getctime) if files else None
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
# Read [ACTIVE FILE] from prompt if present:
active_file = '<path from [ACTIVE FILE: ...]>' if '[ACTIVE FILE:' in PROMPT else latest

TASK TYPE — match ONE, execute ONLY that branch:

[1] LAUNCH APP — task says "open excel application", "open microsoft excel", "launch excel", "start excel"
    AND does NOT mention a specific file to edit:
    os.system('start excel')
    print("EXECUTION_SUCCESS")
    # Done. No file operations.

[2] OPEN FILE — task mentions a specific filename or path (e.g. "open sales.xlsx", "open acm_export.csv"):
    target = 'exact_filename_from_task.xlsx'  # extract from task description (.xlsx or .csv)
    if os.path.isabs(target) and os.path.exists(target):
        found_path = target
    else:
        search_roots = [
            os.path.expanduser('~\\Desktop'),
            os.path.expanduser('~\\OneDrive\\Desktop'),
            os.path.expanduser('~\\Documents'),
            os.path.expanduser('~\\Downloads'),
            'D:\\OneDrive\\Desktop',
            folder,
        ]
        found_path = None
        for root in search_roots:
            if not os.path.exists(root):
                continue
            matches = glob.glob(os.path.join(root, '**', target), recursive=True)
            if matches:
                found_path = matches[0]
                break
    if not found_path:
        raise FileNotFoundError(f"Could not find file: {{target}}")
    os.startfile(found_path)
    print(f"[FILE]: {{found_path}}")
    print("EXECUTION_SUCCESS")
    # Done.

[3] CREATE — task says "create/new/make/generate" + spreadsheet/file:
    save_path = os.path.join(folder, f'descriptive_name_{{ts}}.xlsx')
    wb = Workbook()
    ws = wb.active
    wb.save(save_path)
    print(f"[FILE]: {{save_path}}")
    print("EXECUTION_SUCCESS")
    # Done. Do NOT open the file — keep it unlocked for subsequent edit tasks.

[4] SAVE / CONFIRM — task says "save", "press ok", "press save", "click ok", "click save",
    "confirm saving", "press enter to confirm":
    check_path = active_file
    if check_path and os.path.exists(check_path):
        os.startfile(check_path)  # Open in Excel now that all edits are done
        print(f"[FILE]: {{check_path}}")
    print("EXECUTION_SUCCESS")
    # Done. NO wb.save() call.

[5] DEFAULT — task says write/enter/add/modify/edit/set/insert/update/fill/format:
    wb = load_workbook(active_file) if active_file and os.path.exists(active_file) else Workbook()
    ws = wb.active
    save_path = active_file if active_file and os.path.exists(active_file) else os.path.join(folder, f'spreadsheet_{{ts}}.xlsx')
    # ... apply changes using openpyxl ...
    try:
        wb.save(save_path)
    except PermissionError:
        save_path = os.path.splitext(save_path)[0] + '_v2.xlsx'
        wb.save(save_path)
    print(f"[FILE]: {{save_path}}")
    # Automatically open the file after editing operations are complete
    os.startfile(save_path)
    print("EXECUTION_SUCCESS")

Key openpyxl patterns:
- ws['A1'] = 'value'          # set cell
- ws.cell(row=1, col=1).value = x  # set by row/col
- ws['C1'] = '=A1+B1'         # formula
- wb.save(save_path)

AVOID: subprocess, pyautogui, Excel UI, mouse/keyboard events
"""
    ),
    "powerpoint": ModuleGuidance(
        module_name="PowerPoint",
        library_name="python-pptx",
        library_import="from pptx import Presentation",
        keywords=["powerpoint", "pptx", "slide", "presentation", "bullet point", "layout", "slide show"],
        guidance=f"""
Work with .pptx files using python-pptx. Identify TASK TYPE first, then follow ONLY that branch.

{CRITICAL_DATA_EXTRACTION_RULES}

SETUP:
import os, glob
from pptx import Presentation
from datetime import datetime

# Dynamic folder detection - user's preferred path first
possible_paths = [
    #"D:/OneDrive/Desktop/agent/ppts",  # User's preferred path - FIRST PRIORITY
    os.path.expanduser("~/OneDrive/Desktop/agent/ppts"),
    os.path.expanduser("~/Desktop/agent/ppts"),
    os.path.expanduser("~/Documents/agent/ppts"),
]

folder = None
for path in possible_paths:
    if os.path.exists(os.path.dirname(path)) or os.path.exists(path):
        folder = path
        break

# Final fallback
if not folder:
    folder = os.path.join(os.path.expanduser("~"), "agent", "ppts")

os.makedirs(folder, exist_ok=True)
files = glob.glob(os.path.join(folder, '*.pptx'))
latest = max(files, key=os.path.getctime) if files else None
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
# Read [ACTIVE FILE] from prompt if present:
active_file = '<path from [ACTIVE FILE: ...]>' if '[ACTIVE FILE:' in PROMPT else latest

TASK TYPE — match ONE, execute ONLY that branch:

[1] LAUNCH APP — task says "open powerpoint application", "open microsoft powerpoint", "launch powerpoint", "start powerpoint"
    AND does NOT mention a specific file to edit:
    os.system('start powerpnt')
    print("EXECUTION_SUCCESS")
    # Done. No file operations.

[2] OPEN FILE — task mentions a specific filename or path (e.g. "open slides.pptx", "open existing presentation"):
    target = 'exact_filename_from_task.pptx'  # extract from task description
    if os.path.isabs(target) and os.path.exists(target):
        found_path = target
    else:
        search_roots = [
            os.path.expanduser('~\\Desktop'),
            os.path.expanduser('~\\OneDrive\\Desktop'),
            os.path.expanduser('~\\Documents'),
            os.path.expanduser('~\\Downloads'),
            'D:\\OneDrive\\Desktop',
            folder,
        ]
        found_path = None
        for root in search_roots:
            if not os.path.exists(root):
                continue
            matches = glob.glob(os.path.join(root, '**', target), recursive=True)
            if matches:
                found_path = matches[0]
                break
    if not found_path:
        raise FileNotFoundError(f"Could not find file: {{target}}")
    os.startfile(found_path)
    print(f"[FILE]: {{found_path}}")
    print("EXECUTION_SUCCESS")
    # Done.

[3] CREATE — task says "create/new/make/generate" + presentation/file:
    save_path = os.path.join(folder, f'descriptive_name_{{ts}}.pptx')
    prs = Presentation()
    prs.save(save_path)
    print(f"[FILE]: {{save_path}}")
    print("EXECUTION_SUCCESS")
    # Done. Do NOT open the file — keep it unlocked for subsequent edit tasks.

[4] SAVE / CONFIRM — task says "save", "press ok", "press save", "click ok", "click save",
    "confirm saving", "press enter to confirm":
    check_path = active_file
    if check_path and os.path.exists(check_path):
        os.startfile(check_path)  # Open in PowerPoint now that all edits are done
        print(f"[FILE]: {{check_path}}")
    print("EXECUTION_SUCCESS")
    # Done. NO prs.save() call.

[5] DEFAULT — task says write/type/add/modify/edit/set/insert/format/slide/heading/bullet:
    prs = Presentation(active_file) if active_file and os.path.exists(active_file) else Presentation()
    save_path = active_file if active_file and os.path.exists(active_file) else os.path.join(folder, f'presentation_{{ts}}.pptx')
    # ... apply changes using python-pptx ...
    try:
        prs.save(save_path)
    except PermissionError:
        save_path = os.path.splitext(save_path)[0] + '_v2.pptx'
        prs.save(save_path)
    print(f"[FILE]: {{save_path}}")
    # Automatically open the file after editing operations are complete
    os.startfile(save_path)
    print("EXECUTION_SUCCESS")

Key python-pptx patterns:
- slide = prs.slides.add_slide(prs.slide_layouts[0])   # title slide
- slide = prs.slides.add_slide(prs.slide_layouts[1])   # title + content
- slide.shapes.title.text = 'Title'
- slide.placeholders[1].text = 'Content'
- tf = slide.placeholders[1].text_frame; tf.text = 'bullet'

AVOID: subprocess, pyautogui, PowerPoint UI, mouse/keyboard events
"""
    ),
}


# ============================================================================
# Module Selector
# ============================================================================

class ModuleSelector:
    """
    Intelligently select the best library for a given task.

    Usage:
        selector = ModuleSelector()
        guidance = selector.select_module(task_description)
        if guidance:
            prompt_with_override = selector.inject_override(prompt, guidance)
    """

    def __init__(self):
        """Initialize module selector with predefined modules"""
        self.modules = MODULES
        logger.info(f"[OK] ModuleSelector initialized with {len(self.modules)} modules")

    def select_module(self, task: str) -> Optional[ModuleGuidance]:
        """
        Analyze task and select the most appropriate module.

        Args:
            task: Task description or prompt

        Returns:
            ModuleGuidance if a module matches, None otherwise
        """
        if not task:
            return None

        task_lower = task.lower()

        # Score each module by matching keywords
        scores = {}

        for module_name, guidance in self.modules.items():
            # Count keyword matches
            matches = sum(1 for keyword in guidance.keywords if keyword in task_lower)
            if matches > 0:
                scores[module_name] = matches

        # Return the module with the highest score
        if scores:
            best_module = max(scores, key=scores.get)
            match_count = scores[best_module]
            logger.info(f"[SELECTED] module '{best_module}' ({match_count} keyword matches)")
            return self.modules[best_module]

        logger.debug(f"[NONE] No module selected for task")
        return None

    def inject_override(self, prompt: str, guidance: ModuleGuidance) -> str:
        """
        Inject module override block into the prompt.

        The override is inserted BEFORE the task description so the LLM
        sees it as a hard constraint.

        Args:
            prompt: Original prompt
            guidance: ModuleGuidance object

        Returns:
            Prompt with override block injected at the top
        """
        override_block = guidance.get_override_block()

        # Insert at the beginning, followed by the original prompt
        enriched_prompt = f"{override_block}\n\n{prompt}"

        logger.info(f"[INJECTED] {guidance.module_name} override into prompt")
        return enriched_prompt

    def enrich_prompt(self, prompt: str) -> str:
        """
        Automatically select and inject module override if applicable.

        Args:
            prompt: Original prompt

        Returns:
            Enriched prompt with override (if module matched), or original prompt
        """
        guidance = self.select_module(prompt)

        if guidance:
            return self.inject_override(prompt, guidance)

        return prompt

    def add_custom_module(self, module_name: str, guidance: ModuleGuidance) -> None:
        """
        Add a custom module override at runtime.

        Args:
            module_name: Name of the module (for lookups)
            guidance: ModuleGuidance object
        """
        self.modules[module_name] = guidance
        logger.info(f"[CUSTOM] Added custom module: {module_name}")

    def get_module_info(self, module_name: str) -> Optional[str]:
        """
        Get human-readable info about a module.

        Args:
            module_name: Name of the module

        Returns:
            Formatted string with module info
        """
        if module_name not in self.modules:
            return None

        guidance = self.modules[module_name]
        return f"""
{guidance.module_name} Module
========================================
Library: {guidance.library_name}
Import: {guidance.library_import}

Triggers on keywords:
{', '.join(guidance.keywords)}

Override guidance:
{guidance.guidance}
"""


# ============================================================================
# Quick Helper Functions
# ============================================================================

def select_module(task: str) -> Optional[ModuleGuidance]:
    """Quick function to select a module for a task."""
    selector = ModuleSelector()
    return selector.select_module(task)


def enrich_prompt(prompt: str) -> str:
    """Quick function to enrich a prompt with module overrides."""
    selector = ModuleSelector()
    return selector.enrich_prompt(prompt)


# ============================================================================
# Testing
# ============================================================================

if __name__ == "__main__":
    print("\n" + "="*70)
    print("MODULE SELECTOR TEST")
    print("="*70 + "\n")

    selector = ModuleSelector()

    test_tasks = [
        "Create a Word document with a title and two paragraphs",
        "Create an Excel spreadsheet with sales data",
        "Generate a PowerPoint presentation with 5 slides",
        "Write a Python script using pyautogui",
        "Open file explorer and navigate to Documents",
    ]

    for i, task in enumerate(test_tasks, 1):
        print(f"Test {i}: {task}")
        guidance = selector.select_module(task)

        if guidance:
            print(f"  [YES] Module: {guidance.module_name}")
            print(f"  Library: {guidance.library_name}")
        else:
            print(f"  [NO] No module selected (will use default pyautogui)")

        print()
