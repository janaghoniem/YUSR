from __future__ import annotations

import glob
import os
from datetime import datetime
from typing import Sequence

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.slide import Slide


def _get_agent_folder(subfolder: str) -> str:
    possible_paths = [
        os.path.expanduser("~/OneDrive/Desktop/agent"),
        os.path.expanduser("~/Desktop/agent"),
        os.path.expanduser("~/Documents/agent"),
    ]

    for base_path in possible_paths:
        if os.path.exists(os.path.dirname(base_path)) or os.path.exists(base_path):
            folder = os.path.join(base_path, subfolder)
            os.makedirs(folder, exist_ok=True)
            return folder

    fallback_path = os.path.join(os.path.expanduser("~"), "agent", subfolder)
    os.makedirs(fallback_path, exist_ok=True)
    return fallback_path


def _ensure_pptx_extension(name: str) -> str:
    if name.lower().endswith(".pptx"):
        return name
    return f"{name}.pptx"


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def _print_success(path: str) -> None:
    print(f"[FILE]: {path}")
    print("EXECUTION_SUCCESS")


def _save_with_retry(prs: Presentation, path: str) -> str:
    try:
        prs.save(path)
        return path
    except PermissionError:
        base, ext = os.path.splitext(path)
        retry_path = f"{base}_v2{ext or '.pptx'}"
        prs.save(retry_path)
        return retry_path


def _resolve_ppt_path(filename: str) -> str:
    if os.path.isabs(filename) and os.path.exists(filename):
        return filename

    folder = _get_agent_folder("ppts")
    search_roots = [
        os.path.expanduser("~\\Desktop"),
        os.path.expanduser("~\\OneDrive\\Desktop"),
        os.path.expanduser("~\\Documents"),
        os.path.expanduser("~\\Downloads"),
        folder,
    ]

    for root in search_roots:
        if not os.path.exists(root):
            continue
        direct = os.path.join(root, filename)
        if os.path.isfile(direct):
            return direct
        matches = glob.glob(os.path.join(root, "**", filename), recursive=True)
        if matches:
            return matches[0]

    raise FileNotFoundError(f"Could not find file: {filename}")


def _get_subtitle_placeholder(slide: Slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                return ph
        except Exception:
            continue

    # Common default: subtitle is placeholder index 1 on title slide.
    try:
        return slide.placeholders[1]
    except Exception:
        return None


def _get_body_placeholder(slide: Slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                return ph
        except Exception:
            continue

    # Common default: body/content is placeholder index 1 on title+content slide.
    try:
        return slide.placeholders[1]
    except Exception:
        return None


def ppt_create(name: str | None = None) -> str:
    folder = _get_agent_folder("ppts")
    filename = _ensure_pptx_extension(name) if name else f"presentation_{_timestamp()}.pptx"
    save_path = filename if os.path.isabs(filename) else os.path.join(folder, filename)

    prs = Presentation()
    final_path = _save_with_retry(prs, save_path)
    _print_success(final_path)
    return final_path


def ppt_open(path: str) -> None:
    target = _resolve_ppt_path(path)
    os.startfile(target)
    _print_success(target)


def ppt_find(filename: str) -> str:
    found = _resolve_ppt_path(filename)
    _print_success(found)
    return found


def ppt_load(path: str) -> Presentation:
    resolved = _resolve_ppt_path(path)
    return Presentation(resolved)


def ppt_add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    if slide.shapes.title is not None:
        slide.shapes.title.text = str(title)

    if subtitle:
        subtitle_ph = _get_subtitle_placeholder(slide)
        if subtitle_ph is not None:
            subtitle_ph.text = str(subtitle)

    return prs


def ppt_add_content_slide(prs: Presentation, title: str, content: str) -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    if slide.shapes.title is not None:
        slide.shapes.title.text = str(title)

    body_ph = _get_body_placeholder(slide)
    if body_ph is not None:
        body_ph.text = str(content)

    return prs


def ppt_add_bullet_slide(prs: Presentation, title: str, bullets: Sequence[str]) -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    if slide.shapes.title is not None:
        slide.shapes.title.text = str(title)

    body_ph = _get_body_placeholder(slide)
    if body_ph is None:
        return prs

    if isinstance(bullets, str):
        bullet_items = [b.strip() for b in bullets.splitlines() if b.strip()]
    else:
        bullet_items = [str(b).strip() for b in bullets if str(b).strip()]

    tf = body_ph.text_frame
    tf.clear()

    if not bullet_items:
        return prs

    tf.text = bullet_items[0]
    for item in bullet_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    return prs


def ppt_save(prs: Presentation, path: str) -> str:
    save_path = _ensure_pptx_extension(path)
    if not os.path.isabs(save_path):
        save_path = os.path.join(_get_agent_folder("ppts"), save_path)

    final_path = _save_with_retry(prs, save_path)
    _print_success(final_path)
    return final_path


def ppt_launch() -> None:
    os.system("start powerpnt")
    print("EXECUTION_SUCCESS")


__all__ = [
    "ppt_create",
    "ppt_open",
    "ppt_find",
    "ppt_load",
    "ppt_add_title_slide",
    "ppt_add_content_slide",
    "ppt_add_bullet_slide",
    "ppt_save",
    "ppt_launch",
]
